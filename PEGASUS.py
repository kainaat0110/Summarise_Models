import os
from transformers import PegasusTokenizer, PegasusForConditionalGeneration
import pdfplumber
from pptx import Presentation

# Initialize the Pegasus tokenizer and model
tokenizer = PegasusTokenizer.from_pretrained("google/pegasus-cnn_dailymail")
model = PegasusForConditionalGeneration.from_pretrained("google/pegasus-cnn_dailymail")

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            extracted_text = page.extract_text()
            if extracted_text:  # Check if text extraction was successful
                text += extracted_text + "\n"
    return text

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file."""
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"  # Add newline for separation
    return text

def summarize_text(text):
    """Summarize the provided text using Pegasus model."""
    inputs = tokenizer(text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = model.generate(inputs["input_ids"], max_length=150, min_length=40, length_penalty=2.0, num_beams=4, early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

# Specify your file path
file_path = "/Users/kainaat/Desktop/Summarise Models/Module_2_Part_1.pptx"  # Update with your input file path

# Identify file type and extract text accordingly
file_extension = os.path.splitext(file_path)[1].lower()

if file_extension == ".pdf":
    text = extract_text_from_pdf(file_path)
elif file_extension in [".ppt", ".pptx"]:
    text = extract_text_from_ppt(file_path)
else:
    raise ValueError("Unsupported file format. Please provide a PDF or PPT/PPTX file.")

# Summarize the extracted text
summary = summarize_text(text)

# Print the summary
print("Summary:")
print(summary)
