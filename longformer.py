# Load model directly
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
import pdfplumber
from pptx import Presentation
import os

tokenizer = AutoTokenizer.from_pretrained("Udit191/autotrain-summarization_bart_longformer-54164127153")
model = AutoModelForSeq2SeqLM.from_pretrained("Udit191/autotrain-summarization_bart_longformer-54164127153")

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            extracted_text = page.extract_text()
            if extracted_text:  # Check if text extraction was successful
                text += extracted_text
    return text

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file."""
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"  # Add newline for separation
    return text

def summarize_text(text):
    """Summarize the provided text."""
    inputs = tokenizer.encode("summarize: " + text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = model.generate(inputs, max_length=300, min_length=100, length_penalty=2.0, num_beams=4, early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

# Specify your file path
file_path = "/Users/kainaat/Desktop/Summarise Models/Module_2_Part_1.pptx"  # Update with your input file path

# Identify file type and extract text accordingly
file_extension = os.path.splitext(file_path)[1].lower()

if file_extension == ".pdf":
    text = extract_text_from_pdf(file_path)
elif file_extension in [".ppt", ".pptx"]:
    text = extract_text_from_ppt(file_path)
else:
    raise ValueError("Unsupported file format. Please provide a PDF or PPT/PPTX file.")

# Summarize the extracted text
summary = summarize_text(text)

# Print the summary
print("Summary:")
print(summary)