from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
import pdfplumber
from pptx import Presentation
import os

# Load the Long-T5 model and tokenizer
model_name = "pszemraj/long-t5-tglobal-base-16384-book-summary"
tokenizer = AutoTokenizer.from_pretrained(model_name)
model = AutoModelForSeq2SeqLM.from_pretrained(model_name)

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            extracted_text = page.extract_text()
            if extracted_text:  # Check if text extraction was successful
                text += extracted_text
    return text

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file."""
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"  # Add newline for separation
    return text

def summarize_text(text):
    """Summarize the provided text using Long-T5 model."""
    inputs = tokenizer.encode(text, return_tensors="pt", max_length=4096, truncation=True)  # Adjust max_length as necessary
    summary_ids = model.generate(inputs, max_length=300, min_length=100, length_penalty=2.0, num_beams=4, early_stopping=True)
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

# Specify your file path
file_path = "/Users/kainaat/Desktop/Summarise Models/01_CyberSecurity_LawsStandards.pdf"  # Update with your input file path

# Identify file type and extract text accordingly
file_extension = os.path.splitext(file_path)[1].lower()

if file_extension == ".pdf":
    text = extract_text_from_pdf(file_path)
elif file_extension in [".ppt", ".pptx"]:
    text = extract_text_from_ppt(file_path)
else:
    raise ValueError("Unsupported file format. Please provide a PDF or PPT/PPTX file.")

# Summarize the extracted text
summary = summarize_text(text)

# Print the summary
print("Summary:")
print(summary)
